"""
ASANYX Analytics — Brand Kit Generator

Produces:
  1. /app/public/brand/logos/  — 8 logo variations (PNG)
  2. /app/public/brand/documents/ASANYX_Letterhead.docx
  3. /app/public/brand/documents/ASANYX_Presentation.pptx

Run:  python3 /app/scripts/generate_brand_kit.py
"""
import os
from io import BytesIO
from PIL import Image, ImageDraw, ImageFont, ImageOps, ImageFilter

from docx import Document
from docx.shared import Inches, Pt, RGBColor, Cm, Emu
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches as PI, Pt as PPt, Emu as PE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as PRGBColor

# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------

SRC_LOGO = "/tmp/brandkit/logo_source.png"
OUT_LOGO_DIR = "/app/public/brand/logos"
OUT_DOC_DIR = "/app/public/brand/documents"
os.makedirs(OUT_LOGO_DIR, exist_ok=True)
os.makedirs(OUT_DOC_DIR, exist_ok=True)

BRAND = {
    "deep_blue":   (10, 37, 64),     # #0A2540 — Headings / primary
    "brand_blue":  (30, 107, 255),   # #1E6BFF — CTA / links
    "cyan":        (6, 214, 224),    # #06D6E0 — accent
    "emerald":     (16, 185, 129),   # #10B981 — success
    "white":       (255, 255, 255),
    "soft_grey":   (246, 249, 252),
    "border":      (226, 232, 240),
    "muted":       (100, 116, 139),
    "text_dark":   (15, 23, 42),
}

COMPANY = {
    "name":     "ASANYX ANALYTICS PRIVATE LIMITED",
    "tagline":  "Turning data into decisions that drive growth",
    "website":  "www.asanyxanalytics.com",
    "email":    "hello@asanyxanalytics.com",
    "phone":    "+91 00000 00000",
    "address":  "India · Global Delivery",
}


# ---------------------------------------------------------------------------
# 1. LOGO VARIATIONS
# ---------------------------------------------------------------------------

def _paste_on(bg_color, img, pad_ratio=0.10):
    """Place `img` (RGBA) onto a solid canvas with padding."""
    w, h = img.size
    pad = int(max(w, h) * pad_ratio)
    canvas = Image.new("RGB", (w + pad * 2, h + pad * 2), bg_color)
    canvas.paste(img, (pad, pad), img if img.mode == "RGBA" else None)
    return canvas


def _to_rgba_transparent(rgb_img, white_threshold=245):
    """Convert white background of an RGB logo image to transparent."""
    im = rgb_img.convert("RGBA")
    data = im.getdata()
    new_data = []
    for r, g, b, a in data:
        if r >= white_threshold and g >= white_threshold and b >= white_threshold:
            new_data.append((255, 255, 255, 0))
        else:
            new_data.append((r, g, b, a))
    im.putdata(new_data)
    return im


def _tint(rgba_img, tint_rgb):
    """Recolour all non-transparent pixels to `tint_rgb`, preserving alpha."""
    im = rgba_img.copy()
    px = im.load()
    W, H = im.size
    tr, tg, tb = tint_rgb
    for y in range(H):
        for x in range(W):
            r, g, b, a = px[x, y]
            if a == 0:
                continue
            # Preserve some luminance detail so gradient/shading stays.
            lum = (r + g + b) / (3 * 255)
            # Blend tint with a slight luminance shading — but favour flat tint
            nr = int(tr * (0.60 + 0.40 * lum))
            ng = int(tg * (0.60 + 0.40 * lum))
            nb = int(tb * (0.60 + 0.40 * lum))
            nr = max(0, min(255, nr))
            ng = max(0, min(255, ng))
            nb = max(0, min(255, nb))
            px[x, y] = (nr, ng, nb, a)
    return im


def _crop_mark_from_full(src_rgba):
    """Crop the AS mark (left area) out of the full horizontal logo."""
    W, H = src_rgba.size
    # Empirically: mark sits in ~ x[0.05..0.35] y[0.10..0.72]  (exclude tagline)
    left   = int(W * 0.05)
    top    = int(H * 0.09)
    right  = int(W * 0.34)
    bottom = int(H * 0.72)
    mark = src_rgba.crop((left, top, right, bottom))
    # Trim to bounding box of non-transparent pixels for tight square
    bbox = mark.getbbox()
    if bbox:
        mark = mark.crop(bbox)
    # Make it square with transparent padding
    w, h = mark.size
    side = max(w, h)
    square = Image.new("RGBA", (side, side), (255, 255, 255, 0))
    square.paste(mark, ((side - w) // 2, (side - h) // 2), mark)
    return square


def _crop_horizontal_no_tag(src_rgba):
    """Return horizontal logo (mark + wordmark) without the bottom tagline strip."""
    W, H = src_rgba.size
    return src_rgba.crop((0, 0, W, int(H * 0.78)))


def generate_logos():
    src_rgb = Image.open(SRC_LOGO).convert("RGB")
    src_rgba = _to_rgba_transparent(src_rgb)

    horiz = _crop_horizontal_no_tag(src_rgba)
    # trim to content
    bbox = horiz.getbbox()
    if bbox:
        horiz = horiz.crop(bbox)

    mark = _crop_mark_from_full(src_rgba)

    outputs = []

    # 1. Horizontal color (transparent bg)
    p = os.path.join(OUT_LOGO_DIR, "asanyx-logo-horizontal-color.png")
    horiz.save(p, "PNG"); outputs.append(p)

    # 2. Horizontal on white (safe on light backgrounds)
    on_white = _paste_on(BRAND["white"], horiz, pad_ratio=0.08)
    p = os.path.join(OUT_LOGO_DIR, "asanyx-logo-horizontal-on-white.png")
    on_white.save(p, "PNG"); outputs.append(p)

    # 3. Horizontal on dark navy (for dark backgrounds — tinted white)
    white_horiz = _tint(horiz, BRAND["white"])
    on_dark = _paste_on(BRAND["deep_blue"], white_horiz, pad_ratio=0.10)
    p = os.path.join(OUT_LOGO_DIR, "asanyx-logo-horizontal-on-dark.png")
    on_dark.save(p, "PNG"); outputs.append(p)

    # 4. Horizontal monochrome black (for greyscale printing)
    mono_black = _tint(horiz, BRAND["text_dark"])
    on_white_mono = _paste_on(BRAND["white"], mono_black, pad_ratio=0.08)
    p = os.path.join(OUT_LOGO_DIR, "asanyx-logo-horizontal-mono-black.png")
    on_white_mono.save(p, "PNG"); outputs.append(p)

    # 5. Mark only — color, square, transparent
    p = os.path.join(OUT_LOGO_DIR, "asanyx-mark-color.png")
    mark.save(p, "PNG"); outputs.append(p)

    # 6. Mark on white — social avatar
    m_on_white = _paste_on(BRAND["white"], mark, pad_ratio=0.15)
    p = os.path.join(OUT_LOGO_DIR, "asanyx-mark-on-white.png")
    m_on_white.save(p, "PNG"); outputs.append(p)

    # 7. Mark on deep blue — dark avatar
    mark_white = _tint(mark, BRAND["white"])
    m_on_dark = _paste_on(BRAND["deep_blue"], mark_white, pad_ratio=0.18)
    p = os.path.join(OUT_LOGO_DIR, "asanyx-mark-on-dark.png")
    m_on_dark.save(p, "PNG"); outputs.append(p)

    # 8. Favicon 256×256
    fav = m_on_white.resize((256, 256), Image.LANCZOS)
    p = os.path.join(OUT_LOGO_DIR, "asanyx-favicon-256.png")
    fav.save(p, "PNG"); outputs.append(p)

    return outputs


# ---------------------------------------------------------------------------
# 2. WORD LETTERHEAD
# ---------------------------------------------------------------------------

def _rgbcolor(t):
    return RGBColor(t[0], t[1], t[2])


def _add_horizontal_line(paragraph, color_hex="1E6BFF", size=6):
    """Add a bottom border to a paragraph."""
    pPr = paragraph._p.get_or_add_pPr()
    pBdr = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), str(size))
    bottom.set(qn("w:space"), "1")
    bottom.set(qn("w:color"), color_hex)
    pBdr.append(bottom)
    pPr.append(pBdr)


def generate_letterhead():
    doc = Document()

    # Margins
    for section in doc.sections:
        section.top_margin    = Cm(3.2)
        section.bottom_margin = Cm(2.5)
        section.left_margin   = Cm(2.2)
        section.right_margin  = Cm(2.2)

        # --- HEADER ---
        header = section.header
        htbl = header.add_table(rows=1, cols=2, width=Cm(17))
        htbl.autofit = True
        left_cell, right_cell = htbl.rows[0].cells

        # Left: logo
        logo_path = os.path.join(OUT_LOGO_DIR, "asanyx-logo-horizontal-on-white.png")
        lp = left_cell.paragraphs[0]
        lp.alignment = WD_ALIGN_PARAGRAPH.LEFT
        run = lp.add_run()
        run.add_picture(logo_path, width=Cm(5.5))

        # Right: company block
        rp = right_cell.paragraphs[0]
        rp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        r1 = rp.add_run("ASANYX ANALYTICS")
        r1.font.size = Pt(11); r1.font.bold = True
        r1.font.color.rgb = _rgbcolor(BRAND["deep_blue"])
        rp2 = right_cell.add_paragraph()
        rp2.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        r2 = rp2.add_run("Business Intelligence · Data Engineering · Analytics")
        r2.font.size = Pt(8); r2.font.color.rgb = _rgbcolor(BRAND["muted"])
        rp3 = right_cell.add_paragraph()
        rp3.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        r3 = rp3.add_run(f"{COMPANY['website']}  ·  {COMPANY['email']}")
        r3.font.size = Pt(8); r3.font.color.rgb = _rgbcolor(BRAND["muted"])

        # Header underline
        divider = header.add_paragraph()
        _add_horizontal_line(divider, "1E6BFF", size=8)

        # --- FOOTER ---
        footer = section.footer
        fp = footer.paragraphs[0]
        _add_horizontal_line(fp, "E2E8F0", size=4)
        fp2 = footer.add_paragraph()
        fp2.alignment = WD_ALIGN_PARAGRAPH.CENTER
        fr = fp2.add_run(
            f"{COMPANY['name']}   ·   {COMPANY['address']}   ·   "
            f"{COMPANY['website']}   ·   {COMPANY['email']}"
        )
        fr.font.size = Pt(8); fr.font.color.rgb = _rgbcolor(BRAND["muted"])

    # ---------- BODY: sample letter ----------
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(6)
    r = p.add_run("Date: 22 August 2026")
    r.font.size = Pt(10); r.font.color.rgb = _rgbcolor(BRAND["text_dark"])

    doc.add_paragraph()

    p = doc.add_paragraph()
    r = p.add_run("To,")
    r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = _rgbcolor(BRAND["deep_blue"])

    for line in ["[Recipient Name]", "[Designation]", "[Company Name]", "[Address Line 1]", "[City, State, PIN]"]:
        pp = doc.add_paragraph()
        pr = pp.add_run(line)
        pr.font.size = Pt(10); pr.font.color.rgb = _rgbcolor(BRAND["text_dark"])

    doc.add_paragraph()

    # Subject
    p = doc.add_paragraph()
    r = p.add_run("Subject: ")
    r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = _rgbcolor(BRAND["deep_blue"])
    r2 = p.add_run("[Subject of the letter]")
    r2.font.size = Pt(11); r2.font.color.rgb = _rgbcolor(BRAND["text_dark"])

    doc.add_paragraph()

    body_lines = [
        "Dear [Recipient],",
        "",
        "Thank you for the opportunity to engage with your team. ASANYX Analytics is a boutique "
        "IT consulting practice specialising in Business Intelligence, Data Engineering and modern "
        "cloud analytics on Microsoft Power BI, Fabric, Azure, Databricks and Snowflake.",
        "",
        "This letter serves as a placeholder body. Replace this content with your engagement details, "
        "proposal summary or communication. The letterhead layout, brand colours and footer will "
        "remain consistent across your correspondence.",
        "",
        "We look forward to partnering with you.",
        "",
        "Warm regards,",
    ]
    for line in body_lines:
        pp = doc.add_paragraph()
        run = pp.add_run(line)
        run.font.size = Pt(11); run.font.color.rgb = _rgbcolor(BRAND["text_dark"])

    # Signature block
    p = doc.add_paragraph()
    r = p.add_run("Akshay Birare")
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = _rgbcolor(BRAND["deep_blue"])
    p = doc.add_paragraph()
    r = p.add_run("CEO & Founder  ·  ASANYX Analytics Private Limited")
    r.font.size = Pt(10); r.font.color.rgb = _rgbcolor(BRAND["muted"])

    out_path = os.path.join(OUT_DOC_DIR, "ASANYX_Letterhead.docx")
    doc.save(out_path)
    return out_path


# ---------------------------------------------------------------------------
# 3. POWERPOINT TEMPLATE
# ---------------------------------------------------------------------------

def _add_rect(slide, x, y, w, h, fill_rgb, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRGBColor(*fill_rgb)
    if not line:
        shape.line.fill.background()
    return shape


def _add_text(slide, x, y, w, h, text, size=18, bold=False, color=None, align=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = PPt(size)
    r.font.bold = bold
    if color is not None:
        r.font.color.rgb = PRGBColor(*color)
    return tb


def generate_pptx():
    prs = Presentation()
    prs.slide_width  = PI(13.333)   # 16:9 widescreen
    prs.slide_height = PI(7.5)

    logo_horiz_on_dark = os.path.join(OUT_LOGO_DIR, "asanyx-logo-horizontal-on-dark.png")
    logo_horiz_on_white = os.path.join(OUT_LOGO_DIR, "asanyx-logo-horizontal-on-white.png")
    mark_on_dark = os.path.join(OUT_LOGO_DIR, "asanyx-mark-on-dark.png")

    blank_layout = prs.slide_layouts[6]

    # ------------- SLIDE 1 — Title -------------
    s1 = prs.slides.add_slide(blank_layout)
    # Full dark background
    _add_rect(s1, 0, 0, prs.slide_width, prs.slide_height, BRAND["deep_blue"])
    # Accent bar
    _add_rect(s1, 0, PI(6.9), prs.slide_width, PI(0.12), BRAND["brand_blue"])
    # Logo
    s1.shapes.add_picture(logo_horiz_on_dark, PI(0.6), PI(0.5), height=PI(0.9))
    # Title
    _add_text(s1, PI(0.7), PI(2.6), PI(11.5), PI(1.5),
              "Turning data into decisions that drive growth.",
              size=40, bold=True, color=BRAND["white"])
    # Subtitle
    _add_text(s1, PI(0.7), PI(4.2), PI(11), PI(0.6),
              "Boutique consulting for Business Intelligence, Data Engineering and Cloud Analytics",
              size=18, color=(200, 220, 255))
    # Footer contact
    _add_text(s1, PI(0.7), PI(6.2), PI(12), PI(0.4),
              f"{COMPANY['website']}   ·   {COMPANY['email']}",
              size=12, color=(180, 200, 235))

    # ------------- SLIDE 2 — Section header -------------
    s2 = prs.slides.add_slide(blank_layout)
    _add_rect(s2, 0, 0, prs.slide_width, prs.slide_height, BRAND["soft_grey"])
    _add_rect(s2, 0, 0, PI(0.35), prs.slide_height, BRAND["brand_blue"])
    s2.shapes.add_picture(logo_horiz_on_white, PI(0.7), PI(0.5), height=PI(0.65))
    _add_text(s2, PI(0.7), PI(2.6), PI(12), PI(0.5),
              "01  ·  SECTION",
              size=14, bold=True, color=BRAND["brand_blue"])
    _add_text(s2, PI(0.7), PI(3.1), PI(12), PI(1.6),
              "Section title goes here",
              size=44, bold=True, color=BRAND["deep_blue"])
    _add_text(s2, PI(0.7), PI(4.8), PI(11), PI(1.2),
              "Optional supporting statement for this section — one sentence that frames the "
              "narrative for the slides that follow.",
              size=18, color=BRAND["muted"])

    # ------------- SLIDE 3 — Content (title + bullets) -------------
    s3 = prs.slides.add_slide(blank_layout)
    _add_rect(s3, 0, 0, prs.slide_width, prs.slide_height, BRAND["white"])
    # Header bar
    _add_rect(s3, 0, 0, prs.slide_width, PI(1.0), BRAND["deep_blue"])
    s3.shapes.add_picture(logo_horiz_on_dark, PI(0.5), PI(0.22), height=PI(0.55))
    _add_text(s3, PI(9.5), PI(0.35), PI(3.5), PI(0.35),
              "Business Intelligence", size=11, bold=True,
              color=BRAND["white"])

    _add_text(s3, PI(0.7), PI(1.4), PI(12), PI(1.0),
              "Slide title — one line",
              size=32, bold=True, color=BRAND["deep_blue"])
    _add_text(s3, PI(0.7), PI(2.2), PI(12), PI(0.4),
              "Supporting subtitle or one-liner context.",
              size=14, color=BRAND["muted"])

    # 3 columns
    col_w = PI(3.9); col_h = PI(3.0); top = PI(3.1); gap = PI(0.25)
    left0 = PI(0.7)
    titles = ["Discover", "Design", "Deliver"]
    bodies = [
        "Understand goals, users and current-state data landscape.",
        "Semantic models, dashboards and pipelines built on modern stack.",
        "Migration, rollout and adoption with measurable business impact.",
    ]
    for i, (t, b) in enumerate(zip(titles, bodies)):
        left = left0 + (col_w + gap) * i
        _add_rect(s3, left, top, col_w, col_h, BRAND["soft_grey"])
        _add_rect(s3, left, top, PI(0.08), col_h, BRAND["brand_blue"])
        _add_text(s3, left + PI(0.35), top + PI(0.3), col_w - PI(0.6), PI(0.5),
                  f"0{i+1}", size=12, bold=True, color=BRAND["brand_blue"])
        _add_text(s3, left + PI(0.35), top + PI(0.7), col_w - PI(0.6), PI(0.6),
                  t, size=22, bold=True, color=BRAND["deep_blue"])
        _add_text(s3, left + PI(0.35), top + PI(1.4), col_w - PI(0.6), PI(1.5),
                  b, size=13, color=BRAND["muted"])

    # Footer
    _add_text(s3, PI(0.7), PI(6.9), PI(12), PI(0.3),
              f"ASANYX Analytics   ·   {COMPANY['website']}",
              size=10, color=BRAND["muted"])

    # ------------- SLIDE 4 — Two column narrative -------------
    s4 = prs.slides.add_slide(blank_layout)
    _add_rect(s4, 0, 0, prs.slide_width, prs.slide_height, BRAND["white"])
    _add_rect(s4, 0, 0, prs.slide_width, PI(1.0), BRAND["deep_blue"])
    s4.shapes.add_picture(logo_horiz_on_dark, PI(0.5), PI(0.22), height=PI(0.55))

    _add_text(s4, PI(0.7), PI(1.4), PI(12), PI(1.0),
              "Why ASANYX",
              size=32, bold=True, color=BRAND["deep_blue"])

    _add_text(s4, PI(0.7), PI(2.5), PI(5.7), PI(0.5),
              "PRACTITIONER-LED", size=12, bold=True, color=BRAND["brand_blue"])
    _add_text(s4, PI(0.7), PI(2.9), PI(5.7), PI(3.5),
              "Every engagement is led by senior data & BI practitioners who have "
              "delivered enterprise analytics programs across BFSI, retail, "
              "logistics and edtech — no juniors, no hand-offs.",
              size=15, color=BRAND["text_dark"])

    _add_text(s4, PI(6.9), PI(2.5), PI(5.7), PI(0.5),
              "OUTCOME OBSESSED", size=12, bold=True, color=BRAND["brand_blue"])
    _add_text(s4, PI(6.9), PI(2.9), PI(5.7), PI(3.5),
              "We anchor projects to measurable business outcomes — cost saved, "
              "cycle time cut, revenue unlocked — not just dashboards shipped.",
              size=15, color=BRAND["text_dark"])

    # ------------- SLIDE 5 — Closing / CTA -------------
    s5 = prs.slides.add_slide(blank_layout)
    _add_rect(s5, 0, 0, prs.slide_width, prs.slide_height, BRAND["deep_blue"])
    _add_rect(s5, 0, PI(0), prs.slide_width, PI(0.12), BRAND["brand_blue"])
    s5.shapes.add_picture(mark_on_dark, PI(5.9), PI(1.4), height=PI(1.5))

    _add_text(s5, PI(1), PI(3.2), PI(11.3), PI(1.2),
              "Let's build your modern data platform.",
              size=36, bold=True, color=BRAND["white"])
    _add_text(s5, PI(1), PI(4.5), PI(11.3), PI(0.6),
              "Talk to an ASANYX expert — a boutique team, top-tier craft.",
              size=18, color=(200, 220, 255))

    # CTA "button"
    btn = _add_rect(s5, PI(5.3), PI(5.5), PI(2.7), PI(0.6), BRAND["brand_blue"])
    _add_text(s5, PI(5.3), PI(5.6), PI(2.7), PI(0.5),
              "Talk to an expert →", size=14, bold=True, color=BRAND["white"],
              align=1)  # PP_ALIGN.CENTER = 1

    _add_text(s5, PI(1), PI(6.7), PI(11.3), PI(0.4),
              f"{COMPANY['website']}   ·   {COMPANY['email']}",
              size=12, color=(200, 220, 255))

    out_path = os.path.join(OUT_DOC_DIR, "ASANYX_Presentation.pptx")
    prs.save(out_path)
    return out_path


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    print("[1/3] Generating logo variations …")
    logos = generate_logos()
    for l in logos:
        print("   ✓", l)

    print("[2/3] Generating Word letterhead …")
    doc_path = generate_letterhead()
    print("   ✓", doc_path)

    print("[3/3] Generating PowerPoint template …")
    ppt_path = generate_pptx()
    print("   ✓", ppt_path)

    print("\nDone.")
